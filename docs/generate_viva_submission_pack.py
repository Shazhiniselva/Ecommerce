from pathlib import Path
from datetime import date
import json
from collections import Counter
from typing import Optional

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DIAGRAMS = DOCS / "diagrams"
PRODUCTS_FILE = ROOT / "backend" / "data" / "products.json"
FRONTEND_ASSETS = ROOT / "frontend" / "src" / "assets"

DOCX_OUT = DOCS / "Ecommerce_College_Submission_Report_v3.docx"
PPTX_OUT = DOCS / "Ecommerce_Viva_Ready_Presentation_v3.pptx"

COLLEGE_NAME = "Mohamed Sathak College of Arts & Science"
STUDENT_NAME = "Shalini"
ROLL_NUMBER = "24MCAG001"

# PPT theme colors (formal viva style: white + charcoal + gold)
PPT_BG = PptRGBColor(250, 250, 248)
PPT_SURFACE = PptRGBColor(255, 255, 255)
PPT_PRIMARY = PptRGBColor(30, 30, 34)
PPT_TEXT = PptRGBColor(22, 22, 24)
PPT_MUTED = PptRGBColor(92, 92, 98)
PPT_BORDER = PptRGBColor(214, 201, 171)
PPT_ACCENT = PptRGBColor(173, 132, 66)


def get_font(size=22):
    try:
        return ImageFont.truetype("arial.ttf", size)
    except Exception:
        return ImageFont.load_default()


def draw_box(draw, xy, text, fill=(245, 248, 252), outline=(44, 76, 128)):
    draw.rounded_rectangle(xy, radius=14, fill=fill, outline=outline, width=3)
    x1, y1, x2, y2 = xy
    font = get_font(20)
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text((x1 + (x2 - x1 - tw) / 2, y1 + (y2 - y1 - th) / 2), text, fill=(24, 34, 56), font=font)


def draw_arrow(draw, start, end, color=(55, 55, 55), width=4):
    draw.line([start, end], fill=color, width=width)
    ex, ey = end
    sx, sy = start
    dx = ex - sx
    dy = ey - sy

    if abs(dx) > abs(dy):
        if dx > 0:
            tri = [(ex, ey), (ex - 14, ey - 7), (ex - 14, ey + 7)]
        else:
            tri = [(ex, ey), (ex + 14, ey - 7), (ex + 14, ey + 7)]
    else:
        if dy > 0:
            tri = [(ex, ey), (ex - 7, ey - 14), (ex + 7, ey - 14)]
        else:
            tri = [(ex, ey), (ex - 7, ey + 14), (ex + 7, ey + 14)]
    draw.polygon(tri, fill=color)


def create_backend_module_flow(path: Path):
    img = Image.new("RGB", (1800, 1100), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.text((40, 28), "Backend Module Flow", fill=(18, 50, 92), font=get_font(38))

    draw_box(draw, (100, 220, 470, 360), "Express Server")
    draw_box(draw, (600, 120, 1020, 260), "Routes")
    draw_box(draw, (600, 320, 1020, 460), "Controllers")
    draw_box(draw, (600, 520, 1020, 660), "Middleware")
    draw_box(draw, (600, 720, 1020, 860), "Models")
    draw_box(draw, (1180, 420, 1670, 560), "MongoDB")

    draw_arrow(draw, (470, 290), (600, 190))
    draw_arrow(draw, (470, 290), (600, 390))
    draw_arrow(draw, (470, 290), (600, 590))
    draw_arrow(draw, (470, 290), (600, 790))
    draw_arrow(draw, (1020, 790), (1180, 490))
    draw_arrow(draw, (1020, 390), (1180, 490))

    img.save(path)


def create_auth_flow(path: Path):
    img = Image.new("RGB", (1800, 1200), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.text((40, 28), "Authentication Flow", fill=(18, 50, 92), font=get_font(38))

    nodes = [
        (620, 120, 1180, 230, "User/Admin Login Request"),
        (620, 300, 1180, 410, "Credential Validation"),
        (620, 480, 1180, 590, "JWT Token Issued"),
        (620, 660, 1180, 770, "Protected API Request"),
        (620, 840, 1180, 950, "Middleware Token Verification"),
        (620, 1020, 1180, 1130, "Access Granted / Denied"),
    ]

    for x1, y1, x2, y2, text in nodes:
        draw_box(draw, (x1, y1, x2, y2), text)

    for idx in range(len(nodes) - 1):
        _, _, _, y2, _ = nodes[idx]
        _, ny1, _, _, _ = nodes[idx + 1]
        draw_arrow(draw, (900, y2), (900, ny1))

    img.save(path)


def create_product_lifecycle_flow(path: Path):
    img = Image.new("RGB", (1800, 1100), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.text((40, 28), "Product Lifecycle Flow", fill=(18, 50, 92), font=get_font(38))

    draw_box(draw, (120, 430, 460, 570), "Admin Adds Product")
    draw_box(draw, (560, 210, 980, 350), "Validation + Upload")
    draw_box(draw, (560, 430, 980, 570), "Product Saved in DB")
    draw_box(draw, (560, 650, 980, 790), "List API Response")
    draw_box(draw, (1120, 430, 1660, 570), "Visible in Frontend Collection")

    draw_arrow(draw, (460, 500), (560, 280))
    draw_arrow(draw, (460, 500), (560, 500))
    draw_arrow(draw, (460, 500), (560, 720))
    draw_arrow(draw, (980, 280), (1120, 500))
    draw_arrow(draw, (980, 500), (1120, 500))
    draw_arrow(draw, (980, 720), (1120, 500))

    img.save(path)


def ensure_flowcharts():
    DIAGRAMS.mkdir(parents=True, exist_ok=True)
    create_backend_module_flow(DIAGRAMS / "backend_module_flow.png")
    create_auth_flow(DIAGRAMS / "auth_flow.png")
    create_product_lifecycle_flow(DIAGRAMS / "product_lifecycle_flow.png")


def find_existing_image(candidates) -> Optional[Path]:
    for c in candidates:
        if c.exists():
            return c
    return None


def get_ui_screenshots():
    home = find_existing_image([
        DIAGRAMS / "home_page_screenshot.png",
        DIAGRAMS / "ui_home.png",
        DIAGRAMS / "home_screen.png",
        FRONTEND_ASSETS / "hero_img.png",
    ])
    collection = find_existing_image([
        DIAGRAMS / "collection_page_screenshot.png",
        DIAGRAMS / "ui_collection.png",
        DIAGRAMS / "collection_screen.png",
        FRONTEND_ASSETS / "p_img1.png",
    ])
    return home, collection


def add_page_number(paragraph):
    run = paragraph.add_run()
    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = " PAGE "
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char_begin)
    run._r.append(instr_text)
    run._r.append(fld_char_end)


def heading(doc, text, size=18, color=(16, 52, 96)):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(size)
    r.font.color.rgb = RGBColor(*color)
    return p


def paragraph(doc, text, align=WD_ALIGN_PARAGRAPH.JUSTIFY):
    p = doc.add_paragraph(text)
    p.alignment = align
    p.paragraph_format.space_after = Pt(3)
    p.paragraph_format.line_spacing = 1.15
    return p


def bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(item, style="List Bullet")
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.15


def load_stats():
    if not PRODUCTS_FILE.exists():
        return {
            "count": 0,
            "categories": {},
            "avg": 0,
            "min": 0,
            "max": 0,
        }

    data = json.loads(PRODUCTS_FILE.read_text(encoding="utf-8"))
    prices = [p.get("price", 0) for p in data] or [0]
    categories = Counter(p.get("category", "Unknown") for p in data)
    return {
        "count": len(data),
        "categories": dict(categories),
        "avg": round(sum(prices) / len(prices), 2),
        "min": min(prices),
        "max": max(prices),
    }


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    for row in rows:
        cells = table.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = str(v)


def build_docx():
    stats = load_stats()
    home_shot, collection_shot = get_ui_screenshots()
    doc = Document()

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    # Cover page
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Ecommerce Platform\nProject Submission Report")
    r.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = RGBColor(15, 49, 93)

    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = s.add_run("Academic Submission Edition | Professional Format")
    sr.italic = True
    sr.font.size = Pt(14)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Date: {date.today().isoformat()}")

    doc.add_paragraph()
    doc.add_paragraph()
    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info.add_run(
        f"College: {COLLEGE_NAME}\n"
        f"Student Name: {STUDENT_NAME}\n"
        f"Roll Number: {ROLL_NUMBER}\n"
        "Department: MCA\n"
        "Project Guide: ____________________"
    )

    # Start main section with headers/footers
    main_section = doc.add_section(WD_SECTION.NEW_PAGE)
    header = main_section.header.paragraphs[0]
    header.text = "Ecommerce Platform - Technical Project Report"
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    header.runs[0].font.size = Pt(10)
    header.runs[0].font.color.rgb = RGBColor(90, 90, 90)

    footer = main_section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("Page ")
    add_page_number(footer)

    heading(doc, "Certificate", 20)
    paragraph(doc, "This is to certify that the project report titled 'Ecommerce Platform' is an original work carried out by the candidate under approved academic supervision and is submitted in partial fulfillment of the requirements of the MCA program.")

    doc.add_page_break()
    heading(doc, "Declaration", 20)
    paragraph(doc, "I hereby declare that this report and the software artifacts presented in it are my original work, completed as part of the curriculum requirements under faculty guidance. All external references, frameworks, and learning resources used in this project have been duly acknowledged.")

    doc.add_page_break()
    heading(doc, "Acknowledgement", 20)
    paragraph(doc, "I express my sincere gratitude to my faculty guide, department members, and peers for their support, technical guidance, and continuous feedback throughout the development of this project. Their mentorship played a key role in delivering a complete and stable implementation.")

    doc.add_page_break()
    heading(doc, "Abstract", 20)
    paragraph(doc, "This report presents the design and implementation of a full-stack ecommerce platform composed of a customer storefront, an administrative console, and a backend API service. The solution is implemented using React, Node.js/Express, and MongoDB. The current production-ready baseline supports secure authentication, product catalog management, cart operations, Cash on Delivery (COD) checkout, and end-to-end order lifecycle tracking.")

    doc.add_page_break()
    heading(doc, "Table of Contents", 20)
    index_items = [
        "1. Introduction",
        "2. Objectives",
        "3. Existing System Study",
        "4. Proposed System",
        "5. Architecture Design",
        "6. Technology Stack",
        "7. Frontend Module",
        "8. Admin Module",
        "9. Backend Module",
        "10. Database Design",
        "11. API Design",
        "12. Authentication and Security",
        "13. Checkout and COD Flow",
        "14. Testing Strategy",
        "15. Result Analysis",
        "16. Deployment Guide",
        "17. Viva Questions and Answers",
        "18. Future Enhancements",
        "19. Conclusion",
        "20. References",
        "21. Appendix",
    ]
    bullets(doc, index_items)

    chapters = [
        ("1. Introduction", "The ecommerce platform addresses a complete online retail workflow from product discovery to fulfillment. The project is structured for maintainability, extensibility, and clarity during academic evaluation.", [
            "Supports customer browsing and ordering.",
            "Includes admin controls for products and orders.",
            "Built with real-world full-stack architecture.",
        ]),
        ("2. Objectives", "Project objectives are defined for functionality, stability, and demonstration readiness.", [
            "Build a responsive ecommerce web interface.",
            "Implement secure auth and role-based admin access.",
            "Enable COD-only order flow for operational simplicity.",
            "Provide clean setup documentation for evaluators.",
        ]),
        ("3. Existing System Study", "A review of common ecommerce starter templates identified practical limitations such as weak admin operations, incomplete documentation, and unnecessary payment-gateway complexity for initial releases.", [
            "Lack of clean project modularization.",
            "Unclear deployment setup.",
            "Limited operational reporting structure.",
        ]),
        ("4. Proposed System", "The proposed system introduces a structured three-app setup with clear separation between customer-facing features and administrative operations.", [
            "React storefront for end users.",
            "React admin panel for operations.",
            "Express backend for centralized business logic.",
        ]),
        ("5. Architecture Design", "Architecture follows a client-service-database model. Both frontend clients consume a common backend API, while MongoDB persists entities.", [
            "Low coupling between UI layers.",
            "Reusable middleware and controllers.",
            "Scalable path for future modules.",
        ]),
        ("6. Technology Stack", "The technology stack is selected to balance rapid implementation, maintainability, and production-readiness. The architecture separates frontend, admin, and backend concerns for independent evolution and deployment.", [
            "Frontend Framework: React 18+ for declarative UI rendering with hooks-based state management and component reusability.",
            "Build Tool: Vite for fast HMR (Hot Module Replacement) dev server and optimized production bundling with ES6 module federation.",
            "Styling: Tailwind CSS utility-first framework for consistent design system without custom CSS files.",
            "HTTP Client: Axios for promise-based API requests with request/response interceptors for centralized auth and error handling.",
            "Routing: React Router v6 for client-side navigation with nested route features and lazy code splitting.",
            "Runtime: Node.js 18+ LTS for production-grade stability with native ES modules and improved async performance.",
            "Backend Framework: Express.js minimal but powerful server framework for REST API construction with middleware pipeline architecture.",
            "ODM (Object-Document Mapper): Mongoose for MongoDB schema validation, pre/post hooks, and type safety in document operations.",
            "Authentication: JWT (JSON Web Tokens) with secure token signing/verification for stateless session management across services.",
            "File Handling: Multer middleware for multipart/form-data parsing and file upload handling with stream-based processing.",
            "Database: MongoDB Atlas/Community Edition with replica set support for ACID transactions and horizontal scalability.",
            "Password Security: bcrypt library for irreversible password hashing with configurable work factors against brute-force attacks.",
            "Environment Config: dotenv for .env file parsing enabling environment-specific configurations without hardcoding secrets.",
            "Image Handling: Multer-based upload flow with optional Cloudinary integration for scalable media delivery.",
            "Development Tools: Postman for API endpoint testing and documentation; Git for version control with branching strategy.",
        ]),
        ("7. Frontend Module", "Frontend module handles catalog rendering, product detail navigation, cart updates, and order placement.", [
            "Context-based shared state management.",
            "Responsive page layout and navigation.",
            "Toast notifications for user feedback.",
        ]),
        ("8. Admin Module", "Admin module provides secure login and operational pages for product and order management.", [
            "Add/list/remove products.",
            "View all orders.",
            "Update order status lifecycle.",
        ]),
        ("9. Backend Module", "Backend implementation uses an Express route-controller-model architecture with dedicated middleware for authentication, admin authorization, and file uploads.", [
            "Route groups: /api/user, /api/product, /api/cart, /api/order.",
            "Middleware: auth.js, adminAuth.js, and multer upload pipeline.",
            "Startup pipeline includes DB connect and initial data seeding.",
        ]),
        ("10. Database Design", "MongoDB schemas define users, products, and orders with flexible but controlled document structures.", [
            "Order includes status, paymentMethod, payment flags.",
            "Product includes category, subCategory, size and image arrays.",
            "User schema includes persistent cart object.",
        ]),
        ("11. API Design", "The API follows domain-based route grouping and clear naming conventions to improve maintainability and predictable client integration.", [
            "User APIs: /register, /login, /admin.",
            "Product APIs: /add, /list, /remove, /single.",
            "Cart APIs: /add, /get, /update.",
            "Order APIs: /place, /list, /status, /userorders.",
        ]),
        ("12. Authentication and Security", "Security model combines JWT protection with admin credential checks. Additional hardening recommendations are included from OWASP and Express guidance.", [
            "Secure token verification middleware.",
            "Credential secrets through environment variables.",
            "Recommended enhancements: rate limiting, helmet, audit logging.",
        ]),
        ("13. Checkout and COD Flow", "The checkout pipeline has been simplified to COD-only mode for stable execution and reduced dependency risk.", [
            "Collect address and order lines.",
            "Submit order to backend place endpoint.",
            "Persist order and clear cart on success.",
        ]),
        ("14. Testing Strategy", "Testing combines endpoint verification and user-journey validation to establish confidence before demonstration or deployment.", [
            "API smoke testing for critical routes.",
            "Form validation checks for checkout and login.",
            "Order lifecycle verification from user and admin screens.",
        ]),
        ("15. Result Analysis", f"Dataset-based analysis from available product data provides measurable project evidence. Total products: {stats['count']}, average price: {stats['avg']}, min: {stats['min']}, max: {stats['max']}.", [
            "Category balance helps UI filter behavior.",
            "Price spread supports realistic catalog simulation.",
            "Data quality enables meaningful viva discussion.",
        ]),
        ("16. Deployment Guide", "The deployment process is documented with separate startup commands and environment templates for each application.", [
            "Run backend, frontend, and admin in separate terminals.",
            "Configure environment values from .env.example files.",
            "Seed products for first-time setup.",
        ]),
        ("17. Viva Questions and Answers", "This chapter prepares likely viva questions and concise technical answers for presentation confidence.", [
            "Why COD-only? -> Simpler and stable release path.",
            "How is admin secured? -> Env credentials + JWT token route checks.",
            "How does cart sync work? -> Context state plus backend sync for authenticated users.",
        ]),
        ("18. Future Enhancements", "The platform is extensible and can evolve in phases based on project timeline and scope.", [
            "Online payment reintroduction with secure gateway validation.",
            "Analytics dashboard and sales reporting.",
            "Inventory alerts and coupon engine.",
        ]),
        ("19. Conclusion", "The project demonstrates a complete full-stack ecommerce implementation with modular architecture, practical operational workflows, and professional documentation suitable for academic submission and viva evaluation.", [
            "Functional completeness achieved for core commerce workflow.",
            "Professional documentation and presentation assets prepared.",
            "Ready for academic evaluation and demo.",
        ]),
        ("20. References", "The following references were used to align implementation decisions with widely accepted best practices:", [
            "OWASP API Security Top 10: https://owasp.org/API-Security/",
            "Express Security Best Practices: https://expressjs.com/en/advanced/best-practice-security.html",
            "MongoDB Index Types: https://www.mongodb.com/docs/manual/core/indexes/index-types/",
        ]),
    ]

    for idx, (title_text, overview, points) in enumerate(chapters, start=1):
        doc.add_page_break()
        heading(doc, title_text, 19)
        paragraph(doc, overview)
        bullets(doc, points)

        if idx == 5:
            arch = DIAGRAMS / "architecture_flow.png"
            if arch.exists():
                paragraph(doc, "Figure: High-level architecture flowchart", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(arch), width=Inches(6.2))

        if idx == 7:
            if home_shot is not None:
                paragraph(doc, "Figure: Home page visual", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(home_shot), width=Inches(6.2))
            if collection_shot is not None:
                paragraph(doc, "Figure: Collection page visual", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(collection_shot), width=Inches(6.2))

        if idx == 9:
            backend_flow = DIAGRAMS / "backend_module_flow.png"
            if backend_flow.exists():
                paragraph(doc, "Figure: Backend module flow", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(backend_flow), width=Inches(6.2))

        if idx == 12:
            auth_flow = DIAGRAMS / "auth_flow.png"
            if auth_flow.exists():
                paragraph(doc, "Figure: Authentication flowchart", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(auth_flow), width=Inches(6.2))

        if idx == 13:
            flow = DIAGRAMS / "order_flow.png"
            if flow.exists():
                paragraph(doc, "Figure: COD checkout process flowchart", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(flow), width=Inches(6.2))

        if idx == 15:
            product_flow = DIAGRAMS / "product_lifecycle_flow.png"
            if product_flow.exists():
                paragraph(doc, "Figure: Product lifecycle flowchart", align=WD_ALIGN_PARAGRAPH.LEFT)
                doc.add_picture(str(product_flow), width=Inches(6.2))

    doc.add_page_break()
    heading(doc, "21. Appendix", 19)
    paragraph(doc, "Appendix A: Category Distribution")
    category_rows = sorted([(k, v) for k, v in stats["categories"].items()], key=lambda x: x[0])
    if category_rows:
        add_table(doc, ["Category", "Count"], category_rows)
    else:
        paragraph(doc, "No category records available.")

    doc.add_page_break()
    heading(doc, "Appendix B: Additional Flowcharts", 17)
    appendix_flows = [
        ("Backend Module Flow", DIAGRAMS / "backend_module_flow.png"),
        ("Authentication Flow", DIAGRAMS / "auth_flow.png"),
        ("Product Lifecycle Flow", DIAGRAMS / "product_lifecycle_flow.png"),
    ]
    for caption, path in appendix_flows:
        if path.exists():
            paragraph(doc, f"Figure: {caption}", align=WD_ALIGN_PARAGRAPH.LEFT)
            doc.add_picture(str(path), width=Inches(6.2))

    if home_shot is not None or collection_shot is not None:
        doc.add_page_break()
        heading(doc, "Appendix C: UI Screenshots", 17)
        if home_shot is not None:
            paragraph(doc, "Figure: Home page screenshot", align=WD_ALIGN_PARAGRAPH.LEFT)
            doc.add_picture(str(home_shot), width=Inches(6.2))
        if collection_shot is not None:
            paragraph(doc, "Figure: Collection page screenshot", align=WD_ALIGN_PARAGRAPH.LEFT)
            doc.add_picture(str(collection_shot), width=Inches(6.2))

    DOCS.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX_OUT)


def add_slide_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(34 if len(title) <= 28 else 30)
    p.font.bold = True
    p.font.color.rgb = PPT_PRIMARY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.7))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = PptPt(18)
        sp.font.color.rgb = PPT_MUTED


def add_brand_background(slide):
    # full canvas background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPT_BG
    bg.line.fill.background()

    # top ribbon
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.2))
    top.fill.solid()
    top.fill.fore_color.rgb = PPT_PRIMARY
    top.line.fill.background()

    # footer strip
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3))
    foot.fill.solid()
    foot.fill.fore_color.rgb = PptRGBColor(242, 236, 223)
    foot.line.fill.background()

    # accent underline below title region
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55), Inches(2.3), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PPT_ACCENT
    accent.line.fill.background()


def add_bullet_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide)
    add_slide_title(slide, title, subtitle)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.85), Inches(12.0), Inches(5.0))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PPT_SURFACE
    panel.line.color.rgb = PPT_BORDER

    body = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.8), Inches(4.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = PptPt(8)
    tf.margin_right = PptPt(8)
    tf.margin_top = PptPt(8)
    tf.margin_bottom = PptPt(6)

    bullet_count = len(bullets)
    if bullet_count >= 6:
        body_size = 18
        body_spacing = 4
    elif bullet_count == 5:
        body_size = 20
        body_spacing = 5
    else:
        body_size = 22
        body_spacing = 6

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PptPt(body_size)
        p.font.color.rgb = PPT_TEXT
        p.space_after = PptPt(body_spacing)
        p.line_spacing = 1.12


def add_two_column_slide(prs, title, left_points, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide)
    add_slide_title(slide, title)

    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.9))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = PPT_SURFACE
    left_box.line.color.rgb = PPT_BORDER

    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.9), Inches(5.9), Inches(4.9))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = PPT_SURFACE
    right_box.line.color.rgb = PPT_BORDER

    ltf = left_box.text_frame
    ltf.clear()
    ltf.margin_left = PptPt(10)
    ltf.margin_right = PptPt(10)
    ltf.margin_top = PptPt(8)
    ltf.margin_bottom = PptPt(6)

    max_points = max(len(left_points), len(right_points))
    if max_points >= 6:
        col_size = 16
        col_spacing = 3
    elif max_points == 5:
        col_size = 17
        col_spacing = 4
    else:
        col_size = 19
        col_spacing = 5

    for i, txt in enumerate(left_points):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = txt
        p.font.size = PptPt(col_size)
        p.font.color.rgb = PPT_TEXT
        p.space_after = PptPt(col_spacing)
        p.line_spacing = 1.08

    rtf = right_box.text_frame
    rtf.clear()
    rtf.margin_left = PptPt(10)
    rtf.margin_right = PptPt(10)
    rtf.margin_top = PptPt(8)
    rtf.margin_bottom = PptPt(6)
    for i, txt in enumerate(right_points):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = txt
        p.font.size = PptPt(col_size)
        p.font.color.rgb = PPT_TEXT
        p.space_after = PptPt(col_spacing)
        p.line_spacing = 1.08


def add_image_slide(prs, title, image_path, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide)
    add_slide_title(slide, title)
    if image_path.exists():
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.65), Inches(11.4), Inches(5.15))
        panel.fill.solid()
        panel.fill.fore_color.rgb = PPT_SURFACE
        panel.line.color.rgb = PPT_BORDER
        slide.shapes.add_picture(str(image_path), Inches(1.2), Inches(1.8), width=Inches(10.9))
        cap = slide.shapes.add_textbox(Inches(1.2), Inches(6.6), Inches(10.9), Inches(0.5))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        p.font.size = PptPt(14)
        p.font.color.rgb = PPT_MUTED


def build_pptx():
    stats = load_stats()
    home_shot, collection_shot = get_ui_screenshots()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title Slide
    title = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(title)
    add_slide_title(title, "Ecommerce Platform", "Academic Viva Presentation | Professional Edition")

    badge = title.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(6.0), Inches(1.0))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PPT_ACCENT
    badge.line.fill.background()
    btf = badge.text_frame
    bp = btf.paragraphs[0]
    bp.text = "College Project Demonstration"
    bp.font.color.rgb = PptRGBColor(255, 255, 255)
    bp.font.bold = True
    bp.font.size = PptPt(22)

    details = title.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.8), Inches(2.0))
    details_panel = title.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.35), Inches(11.95), Inches(2.15))
    details_panel.fill.solid()
    details_panel.fill.fore_color.rgb = PPT_SURFACE
    details_panel.line.color.rgb = PPT_BORDER
    details = title.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.6), Inches(1.8))
    dtf = details.text_frame
    d1 = dtf.paragraphs[0]
    d1.text = f"College: {COLLEGE_NAME}"
    d1.font.size = PptPt(20)
    d1.font.color.rgb = PPT_TEXT

    d2 = dtf.add_paragraph()
    d2.text = f"Student: {STUDENT_NAME} | Roll No: {ROLL_NUMBER}"
    d2.font.size = PptPt(20)
    d2.font.color.rgb = PPT_TEXT

    d3 = dtf.add_paragraph()
    d3.text = f"Department: MCA"
    d3.font.size = PptPt(18)
    d3.font.color.rgb = PPT_MUTED

    # SLIDE 2: Presentation Outline
    add_bullet_slide(prs, "Presentation Outline", [
        "1. Problem Statement & Objectives",
        "2. Technology Architecture & Design",
        "3. Frontend Implementation Details",
        "4. Backend & API Design",
        "5. Database Schema & Models",
        "6. Security & Authentication",
        "7. Checkout & COD Flow",
        "8. Testing, Results & Deployment"
    ])

    # SLIDE 3: Problem & Objectives
    add_bullet_slide(prs, "Problem Statement", [
        "Learning Need: A complete, production-style full-stack implementation for academic evaluation",
        "Architecture Need: Clear separation of concerns across frontend, admin, and backend",
        "Operational Need: Stable COD-first checkout for reliable demonstrations",
        "Documentation Need: Structured technical report and viva-ready presentation",
        "Scalability Need: Foundation that can evolve toward payment gateways and analytics"
    ])

    # SLIDE 4: Project Objectives
    add_two_column_slide(prs, "Project Objectives",
        ["Functional Goals:", "• Build responsive customer interface", "• Implement admin panel", "• Design secure backend APIs", "• Enable complete order pipeline"],
        ["Quality Goals:", "• Professional architecture", "• Comprehensive documentation", "• Security best practices", "• Ready for deployment"]
    )

    # SLIDE 5: Technology Stack - Part 1
    add_two_column_slide(prs, "Frontend Technology Stack",
        ["Framework & Build:", "• React 18+ (Hooks-based)", "• Vite (Fast HMR/bundling)", "• Tailwind CSS (Utility-first)", "• React Router v6 (Client navigation)"],
        ["Libraries & Tools:", "• Axios (HTTP client)", "• Context API (State mgmt)", "• React Toastify (Notifications)", "• npm/yarn (Package manager)"]
    )

    # SLIDE 6: Technology Stack - Part 2
    add_two_column_slide(prs, "Backend & Database Stack",
        ["Backend Runtime:", "• Node.js 18+ LTS", "• Express.js (Minimal framework)", "• Mongoose ODM (Schema validation)", "• JWT (Token-based auth)"],
        ["Infrastructure:", "• MongoDB (Document database)", "• Multer (File uploads)", "• bcrypt (Password hashing)", "• Environment variables (.env)"]
    )

    # SLIDE 7: Architecture Overview
    add_bullet_slide(prs, "System Architecture", [
        "Three-Tier Separation: Customer Frontend | Admin Panel | Shared Backend",
        "API-First Design: Both clients communicate via REST endpoints",
        "Stateless Authentication: JWT tokens for session management",
        "Database Layer: MongoDB with schema validation via Mongoose",
        "Scalability Path: Ready for microservices, caching, and load balancing"
    ])

    # SLIDE 8: Architecture Diagram
    add_image_slide(prs, "System Architecture Flowchart", DIAGRAMS / "architecture_flow.png", 
        "Client-Service-Database model: Frontend and Admin clients consume unified backend API")

    # SLIDE 9: Frontend Module
    add_bullet_slide(prs, "Frontend Module Features", [
        "Product Discovery: Search, filter by category, price range sorting",
        "Product Details: Full specs, images, size/color variants",
        "Shopping Cart: Add/remove, quantity management, persistent storage",
        "Checkout: Address entry, order review, COD confirmation",
        "Order Tracking: History view, status updates, delivery confirmation"
    ])

    # SLIDE 10: Admin Module
    add_two_column_slide(prs, "Admin Module Capabilities",
        ["Product Management:", "• Add new products", "• Edit product details", "• Remove/archive", "• Bulk operations"],
        ["Order Management:", "• View all orders", "• Update status", "• Track fulfillment", "• Generate reports"]
    )

    # SLIDE 11: Backend API Structure
    add_bullet_slide(prs, "Backend API Organization", [
        "User Routes: /api/user (register, login, admin login)",
        "Product Routes: /api/product (add, list, remove, single)",
        "Cart Routes: /api/cart (add, get, update)",
        "Order Routes: /api/order (place, list, status, userorders)",
        "Route-level middleware enforces auth and admin access controls"
    ])

    # SLIDE 12A: Backend Flow
    add_image_slide(prs, "Backend Module Flow", DIAGRAMS / "backend_module_flow.png",
        "Express route -> controller -> middleware -> model pipeline with MongoDB persistence")

    # SLIDE 13: Database Schema
    add_two_column_slide(prs, "MongoDB Collections",
        ["User Schema:", "• ID, Email, Password Hash", "• Name, Address, Phone", "• Cart Array (Items)", "• Auth Tokens, Timestamps"],
        ["Product Schema:", "• ID, Name, Category", "• Price, Description", "• Images Array", "• Sizes, Stock Level"]
    )

    # SLIDE 14: Order Model
    add_two_column_slide(prs, "Order & Payment Models",
        ["Order Schema:", "• User Reference", "• Items Array", "• Delivery Address", "• Order Date & Status"],
        ["Payment Flow:", "• Payment Method: COD", "• Status Transitions", "• Date Tracking", "• Admin Updates"]
    )

    # SLIDE 15: Authentication
    add_bullet_slide(prs, "Authentication & Security", [
        "JWT Implementation: Secure token generation and verification",
        "Password Security: bcrypt hashing with salt rounds",
        "Role-Based Access: Customer vs Admin vs Public routes",
        "Middleware Protection: Auth middleware on protected endpoints",
        "Environment Secrets: Credentials via .env files, never hardcoded"
    ])

    # SLIDE 16: Auth Flowchart
    add_image_slide(prs, "Authentication Flowchart", DIAGRAMS / "auth_flow.png",
        "Credential validation, JWT issuance, and middleware-based access control")

    # SLIDE 17: COD Checkout Flow
    add_bullet_slide(prs, "Cash On Delivery Flow", [
        "Step 1: Customer reviews cart items and prices",
        "Step 2: Enter delivery address and phone number",
        "Step 3: Confirm order submission (no payment gateway call)",
        "Step 4: Backend validates and creates order document",
        "Step 5: Cart cleared, order appears in history and admin dashboard"
    ])

    # SLIDE 18: Checkout Process Diagram
    add_image_slide(prs, "Order Processing Pipeline", DIAGRAMS / "order_flow.png", 
        "COD checkout simplified: No external gateway, single backend endpoint")

    # SLIDE 19: Admin Operations
    add_image_slide(prs, "Admin Workflow", DIAGRAMS / "admin_flow.png", 
        "Admin authentication and operational monitoring lifecycle")

    # SLIDE 20: Product Lifecycle Diagram
    add_image_slide(prs, "Product Lifecycle Flowchart", DIAGRAMS / "product_lifecycle_flow.png",
        "Admin add/update actions reflected in catalog API responses and storefront rendering")

    # SLIDE 21: UI Screenshots
    if home_shot is not None:
        add_image_slide(prs, "Application UI - Home", home_shot,
            "Home page layout used in the customer storefront")
    if collection_shot is not None:
        add_image_slide(prs, "Application UI - Collection", collection_shot,
            "Collection and filtering interface used in product browsing")

    # SLIDE 18: Data Statistics
    add_bullet_slide(prs, "Dataset & Metrics", [
        f"Total Products: {stats['count']}",
        f"Price Range: ₹{stats['min']} - ₹{stats['max']}",
        f"Average Price: ₹{stats['avg']}",
        "Categories: " + ", ".join(list(stats['categories'].keys())[:5]),
        "Ready for realistic catalog simulation and viva discussion"
    ])

    # SLIDE 19: Testing Strategy
    add_two_column_slide(prs, "Testing & Validation",
        ["Unit Testing:", "• Form validation logic", "• Cart calculations", "• Auth token checks", "• API error handling"],
        ["Integration Testing:", "• API smoke tests", "• User registration flow", "• Product listing", "• Order placement"]
    )

    # SLIDE 20: Security Best Practices
    add_bullet_slide(prs, "Security Implementation", [
        "HTTPS/TLS: Encrypted data transmission",
        "CORS Configuration: Controlled cross-origin requests",
        "Input Validation: Sanitize all user inputs",
        "NoSQL Injection Mitigation: Validate and sanitize query/body payloads",
        "Recommended Additions: Helmet.js, rate limiting, audit logging"
    ])

    # SLIDE 21: Deployment Strategy
    add_bullet_slide(prs, "Deployment & Scalability", [
        "Local Development: Separate npm processes (backend:4000, frontend/admin via Vite)",
        "Environment Configuration: .env files per deployment stage",
        "Database: MongoDB Atlas cloud or local community edition",
        "Future Scale: Containerization (Docker), CI/CD pipeline, load balancing",
        "Uptime: Monitoring and alerting infrastructure"
    ])

    # SLIDE 22: Key Features Summary
    add_two_column_slide(prs, "Project Highlights",
        ["Complete Features:", "✓ User authentication", "✓ Product catalog", "✓ Shopping cart", "✓ Order lifecycle"],
        ["Quality Marks:", "✓ Clean code structure", "✓ Security hardened", "✓ Well documented", "✓ Production-ready"]
    )

    # SLIDE 23: Viva Q&A - Why COD?
    add_bullet_slide(prs, "Viva Q&A: Why COD-Only?", [
        "Simplicity: Single payment pathway reduces complexity",
        "Stability: No dependency on third-party payment gateway",
        "MVP Approach: Focuses on core ecommerce logic, not payment integration",
        "Academic Clarity: Demonstrates order pipeline without payment noise",
        "Future Ready: Payment gateways can be reintroduced in Phase 2"
    ])

    # SLIDE 24: Viva Q&A - Security
    add_bullet_slide(prs, "Viva Q&A: Security Model", [
        "Admin Credentials: Stored as environment variables, not in code",
        "JWT Tokens: Issued on login, verified on protected routes",
        "Password Storage: bcrypt hashing prevents plaintext exposure",
        "CORS Policy: Controls which domains can access our API",
        "Rate Limiting: (Recommended) Prevents brute-force attacks"
    ])

    # SLIDE 25: Viva Q&A - Cart Sync
    add_bullet_slide(prs, "Viva Q&A: Cart Synchronization", [
        "Client State: Frontend context updates cart instantly for responsive UX",
        "Server Persistence: Cart data is synchronized through /api/cart endpoints",
        "Authentication: Token-based calls protect user-specific cart data",
        "Consistency: Backend remains the source of truth for order creation",
        "Reliability: Sync after mutations prevents stale cart snapshots"
    ])

    # SLIDE 26: Viva Q&A - Scalability
    add_bullet_slide(prs, "Viva Q&A: Scaling Approach", [
        "Database: MongoDB replica sets for read scaling",
        "Caching: Redis layer for product catalog and hot data",
        "Microservices: Split into payment, inventory, notification services",
        "Load Balancing: Nginx/HAProxy across multiple node instances",
        "Monitoring: Prometheus + Grafana for performance tracking"
    ])

    # SLIDE 27: Viva Q&A - Challenges & Solutions
    add_two_column_slide(prs, "Challenges Addressed",
        ["Technical:", "• State management", "• Real-time updates", "• Data consistency", "• Error handling"],
        ["Solutions:", "• Context API for state", "• WebSocket readiness", "• Transaction logs", "• Graceful degradation"]
    )

    # SLIDE 28: Project Timeline
    add_bullet_slide(prs, "Development Timeline", [
        "Phase 1 (Backend): Database design, API development, testing",
        "Phase 2 (Frontend): UI components, integration, cart logic",
        "Phase 3 (Admin): Dashboard, product management, order tracking",
        "Phase 4 (Documentation): API docs, setup guide, viva materials",
        "Phase 5 (Deployment): Testing, security audit, launch readiness"
    ])

    # SLIDE 29: Future Enhancements
    add_two_column_slide(prs, "Roadmap & Enhancements",
        ["Short-term:", "• Payment gateway integration", "• Inventory alerts", "• Email notifications", "• User reviews"],
        ["Long-term:", "• Analytics dashboard", "• Coupon system", "• Recommendation engine", "• Mobile app"]
    )

    # SLIDE 30: Conclusion
    add_bullet_slide(prs, "Project Conclusion", [
        "Complete full-stack ecommerce solution with a stable COD-first workflow",
        "Professional documentation suitable for institutional evaluation",
        "Security-aware implementation aligned with OWASP recommendations",
        "Modular architecture ready for iterative enhancements",
        "Prepared for viva presentation and guided deployment"
    ])

    # SLIDE 31: Thank You
    thank_you = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(thank_you)
    
    ty_title = thank_you.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    ty_tf = ty_title.text_frame
    ty_p = ty_tf.paragraphs[0]
    ty_p.text = "Thank You"
    ty_p.font.size = PptPt(54)
    ty_p.font.bold = True
    ty_p.font.color.rgb = PPT_PRIMARY
    ty_p.alignment = PP_ALIGN.CENTER

    ty_sub = thank_you.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(1.5))
    ty_stf = ty_sub.text_frame
    ty_sp = ty_stf.paragraphs[0]
    ty_sp.text = f"Questions & Discussion\n{STUDENT_NAME} | {COLLEGE_NAME}"
    ty_sp.font.size = PptPt(24)
    ty_sp.font.color.rgb = PptRGBColor(90, 90, 90)
    ty_sp.alignment = PP_ALIGN.CENTER

    DOCS.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))


def main():
    ensure_flowcharts()
    build_docx()
    build_pptx()
    print(f"Created: {DOCX_OUT}")
    print(f"Created: {PPTX_OUT}")


if __name__ == "__main__":
    main()
