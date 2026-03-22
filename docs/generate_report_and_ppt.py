from pathlib import Path
from datetime import date

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION_START

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs"
DOCX_PATH = OUT_DIR / "Ecommerce_Project_Report_30_Pages.docx"
PPTX_PATH = OUT_DIR / "Ecommerce_Project_Presentation.pptx"


def add_header(paragraph, text, size=18):
    run = paragraph.add_run(text)
    run.bold = True
    run.font.size = Pt(size)
    run.font.name = "Calibri"


def build_docx():
    doc = Document()

    # Global style
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    # Title page
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("Ecommerce Platform\nTechnical and Business Report")
    r.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = RGBColor(31, 56, 100)

    doc.add_paragraph()
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rs = subtitle.add_run("Professional Project Documentation")
    rs.italic = True
    rs.font.size = Pt(15)

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_run = meta.add_run(f"Prepared on {date.today().isoformat()}")
    meta_run.font.size = Pt(12)

    doc.add_page_break()

    sections = [
        ("1. Executive Summary", [
            "This report presents the architecture, implementation, and operational strategy of a full-stack ecommerce platform with customer storefront, admin panel, and backend services.",
            "The platform currently focuses on COD-only payment flow to simplify operations and reduce payment gateway dependency risks.",
            "This document explains technical design, deployment approach, security controls, testing strategy, and growth roadmap."
        ]),
        ("2. Project Objectives", [
            "Deliver a responsive customer-facing ecommerce experience for browsing, cart management, and ordering.",
            "Provide an admin dashboard for product lifecycle and order operations.",
            "Ensure maintainable code structure with clear separation of frontend, admin frontend, and backend APIs.",
            "Support local development and production deployment with predictable setup steps."
        ]),
        ("3. Scope and Deliverables", [
            "Frontend application for customers built with React and Vite.",
            "Admin panel application built with React and Vite.",
            "Node.js and Express backend with MongoDB persistence.",
            "Seed scripts, environment templates, and setup documentation.",
            "COD-only order processing and status tracking workflows."
        ]),
        ("4. Technology Stack", [
            "Frontend and Admin: React, React Router, Axios, Toastify, Tailwind CSS, Vite.",
            "Backend: Node.js, Express, Mongoose, JWT auth, Multer for image uploads.",
            "Database: MongoDB local instance (default), with configurable DB name.",
            "Tooling: npm, ESLint, dotenv-based runtime configuration."
        ]),
        ("5. System Architecture", [
            "The solution follows a multi-application architecture with shared backend APIs.",
            "Customer and admin interfaces are independent frontends that consume backend endpoints.",
            "Domain logic is organized into controllers, routes, middleware, and models.",
            "This separation enables independent UI evolution while preserving centralized business rules."
        ]),
        ("6. Repository Structure", [
            "Root workspace contains frontend, backend, and admin applications.",
            "Backend includes configs, controllers, middleware, models, routes, and scripts.",
            "Frontend and admin include components, pages, assets, and entry points.",
            "This structure supports team parallelism and modular maintenance."
        ]),
        ("7. Environment Configuration", [
            "Environment variables are documented through .env.example files for all apps.",
            "Backend requires PORT, JWT_SECRET, ADMIN_EMAIL, ADMIN_PASSWORD, and MongoDB URI settings.",
            "Frontend and admin rely on VITE_BACKEND_URL to connect to backend APIs.",
            "Configuration-driven setup reduces hardcoded values and improves deployment portability."
        ]),
        ("8. Backend API Design", [
            "REST endpoints are grouped by user, product, cart, and order domains.",
            "Controllers hold business logic while routes map HTTP verbs and paths.",
            "JWT-based middleware secures protected operations.",
            "Consistent JSON responses improve frontend integration reliability."
        ]),
        ("9. Authentication and Authorization", [
            "User authentication uses credential validation and JWT token issuance.",
            "Admin authentication is controlled by backend environment credentials.",
            "Protected routes verify tokens before action execution.",
            "Role-specific access paths reduce accidental privilege escalation."
        ]),
        ("10. Data Model Overview", [
            "Core models include users, products, and orders.",
            "Order schema records items, address, payment method, payment status, and lifecycle status.",
            "Product schema supports category and size-based filtering use cases.",
            "Mongoose models provide schema-level consistency and query ergonomics."
        ]),
        ("11. Product Management Flow", [
            "Admin users add products through controlled forms and backend endpoints.",
            "Image assets are managed through local upload paths in current setup.",
            "Product listing and removal are available through admin operations.",
            "Frontend consumes product list endpoint for catalog rendering."
        ]),
        ("12. Catalog and Collection Experience", [
            "Frontend collection page renders products from backend list endpoint.",
            "Search and category logic improve product discoverability.",
            "Product detail pages provide size selections and cart actions.",
            "Seeded product dataset ensures immediate usability in local environments."
        ]),
        ("13. Cart Management", [
            "Cart state is managed in React context and synchronized with backend for authenticated users.",
            "Quantity and size combinations are tracked with nested data structures.",
            "Cart totals are computed based on product price and quantity.",
            "Robust cart flows are foundational for order conversion."
        ]),
        ("14. Order Placement Workflow", [
            "Checkout captures delivery information and validates required fields.",
            "Order payload includes address, selected items, and computed total.",
            "Backend stores order and clears cart when placement succeeds.",
            "Frontend redirects users to order history after successful submission."
        ]),
        ("15. Payment Strategy (COD Only)", [
            "The system currently supports Cash On Delivery as the only payment method.",
            "Stripe and Razorpay flows were removed from frontend and backend for operational simplification.",
            "Order records still preserve payment method metadata for reporting.",
            "This approach minimizes setup friction for local and initial production rollout."
        ]),
        ("16. Admin Panel Operations", [
            "Admin panel supports login, product management, and order tracking.",
            "Order status updates are exposed via protected backend endpoints.",
            "Operational visibility helps teams process and fulfill orders efficiently.",
            "Administrative UX is optimized for routine commerce operations."
        ]),
        ("17. Order Lifecycle Management", [
            "Orders transition through statuses managed by admin actions.",
            "Users can view personal order history in the storefront.",
            "Status visibility aligns customer expectations during fulfillment.",
            "Lifecycle data enables future analytics and SLA tracking."
        ]),
        ("18. Error Handling and User Feedback", [
            "Frontend uses toast notifications for user-facing success and failure messages.",
            "Backend returns structured error messages for operational clarity.",
            "Try-catch patterns around async operations reduce uncaught runtime failures.",
            "Clear feedback loops improve trust and usability."
        ]),
        ("19. Database Strategy", [
            "MongoDB local fallback is configured for seamless local development.",
            "Database name is configurable through environment variables.",
            "Seeder script initializes product catalog quickly for demonstrations and QA.",
            "Model-driven access patterns support future migration and scaling."
        ]),
        ("20. Local Development Setup", [
            "Each app is started in its own terminal process.",
            "Backend default runs on port 4000; frontend and admin run on Vite ports.",
            "Setup documents and env templates reduce onboarding time.",
            "Developers can validate end-to-end flow in minutes."
        ]),
        ("21. Deployment Readiness", [
            "Backend supports configurable ports and environment-driven secrets.",
            "Frontend/admin build outputs can be deployed independently.",
            "CORS and API URL configuration are already integrated.",
            "Additional observability and CI checks are recommended before large-scale launch."
        ]),
        ("22. Security Considerations", [
            "JWT secrets and admin credentials must be stored securely in environment variables.",
            "Input validation should be enforced consistently across all mutable endpoints.",
            "Transport security (HTTPS) is mandatory in production.",
            "Rate limiting and audit logging are recommended next controls."
        ]),
        ("23. Performance Considerations", [
            "Catalog APIs should support pagination for large product catalogs.",
            "Client-side rendering can be optimized with lazy loading strategies.",
            "Image optimization and CDN strategies will reduce page load times.",
            "Database indexing should be applied for high-volume query paths."
        ]),
        ("24. Testing Strategy", [
            "Unit tests should cover controllers, validators, and utility layers.",
            "Integration tests should validate auth, cart, and order APIs.",
            "UI regression tests are recommended for checkout and admin order updates.",
            "A release checklist should gate deployment quality."
        ]),
        ("25. Operational Risks and Mitigations", [
            "Risk: Missing environment config can break startup. Mitigation: enforce .env validation at boot.",
            "Risk: Port conflicts in local runs. Mitigation: document dedicated ports and health checks.",
            "Risk: Data inconsistency in manual admin actions. Mitigation: add stronger status transition rules.",
            "Risk: Scaling bottlenecks. Mitigation: monitor and optimize hotspots iteratively."
        ]),
        ("26. UX and Accessibility Notes", [
            "Responsive layouts support common desktop and mobile breakpoints.",
            "Form fields use labels and clear placeholders to reduce user error.",
            "Feedback patterns improve navigation confidence in key flows.",
            "Accessibility audits should be part of future release criteria."
        ]),
        ("27. Change Log Highlights", [
            "Migrated to COD-only payment method.",
            "Removed Stripe and Razorpay integration points.",
            "Removed cloud payment dependencies in backend order flow.",
            "Added setup documentation and environment templates."
        ]),
        ("28. Future Enhancements", [
            "Role-based admin users and permission granularity.",
            "Coupon and promotions engine for conversion optimization.",
            "Inventory synchronization and low-stock alerts.",
            "Analytics dashboard for sales and customer behavior insights."
        ]),
        ("29. Conclusion", [
            "The ecommerce platform is operationally complete for core browsing, ordering, and admin management.",
            "Current COD-first strategy improves implementation stability and operational simplicity.",
            "With targeted enhancements, the system can scale into a production-grade commerce platform.",
            "This report provides a professional baseline for engineering, business, and stakeholder alignment."
        ]),
        ("30. Appendix", [
            "Important directories: backend, frontend, admin.",
            "Core commands: npm install, npm run dev, npm run server, node seeder.js.",
            "Admin auth endpoint: POST /api/user/admin.",
            "Default local DB URI: mongodb://127.0.0.1:27017/e-commerce."
        ]),
    ]

    for idx, (heading, points) in enumerate(sections, start=1):
        h = doc.add_paragraph()
        add_header(h, heading, size=18)
        h.alignment = WD_ALIGN_PARAGRAPH.LEFT

        doc.add_paragraph()
        for p in points:
            para = doc.add_paragraph(p, style="List Bullet")
            para.paragraph_format.space_after = Pt(8)

        note = doc.add_paragraph()
        nr = note.add_run("Key takeaway: this section strengthens implementation clarity and execution confidence.")
        nr.italic = True
        nr.font.color.rgb = RGBColor(80, 80, 80)

        if idx < len(sections):
            doc.add_page_break()

    doc.save(DOCX_PATH)


def add_title_only_slide(prs, title_text, subtitle_text=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    title.text_frame.paragraphs[0].font.size = PptPt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(25, 62, 109)

    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.size = PptPt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(70, 70, 70)


def add_bullets_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(22)
        p.font.name = "Calibri"
        p.font.color.rgb = PptRGBColor(40, 40, 40)


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    add_title_only_slide(
        prs,
        "Ecommerce Platform",
        "Professional Project Presentation | COD-Only Release"
    )

    slides = [
        ("Agenda", [
            "Project context and goals",
            "Architecture and technology",
            "Checkout and admin operations",
            "Deployment readiness and roadmap"
        ]),
        ("Project Vision", [
            "Create a reliable ecommerce experience",
            "Enable efficient admin operations",
            "Use maintainable full-stack architecture",
            "Prepare for scalable growth"
        ]),
        ("Business Objectives", [
            "Improve product discovery and ordering",
            "Reduce operational complexity with COD-only flow",
            "Accelerate team onboarding with clear setup docs",
            "Support future payment and analytics expansion"
        ]),
        ("Solution Overview", [
            "Customer frontend (React + Vite)",
            "Admin frontend (React + Vite)",
            "Backend API (Node.js + Express + MongoDB)",
            "Shared operational data and auth controls"
        ]),
        ("Architecture", [
            "Frontend and admin consume centralized backend APIs",
            "Controllers + routes + middleware pattern",
            "Mongoose models enforce data consistency",
            "Environment-driven runtime configuration"
        ]),
        ("Technology Stack", [
            "React, React Router, Axios, Tailwind",
            "Node.js, Express, Mongoose, JWT",
            "MongoDB local default for development",
            "Vite build pipeline for modern development"
        ]),
        ("Core User Flows", [
            "Browse products and collections",
            "Add items to cart with size variants",
            "Place COD order with address details",
            "Track order history and status"
        ]),
        ("Admin Flows", [
            "Admin login via backend credentials",
            "Create and manage products",
            "View complete order list",
            "Update order statuses"
        ]),
        ("Payment Simplification", [
            "Stripe and Razorpay removed",
            "Cash On Delivery is the only method",
            "Lower setup and maintenance overhead",
            "Faster issue diagnosis and onboarding"
        ]),
        ("Data and Models", [
            "User, Product, and Order schemas",
            "Order status + payment method metadata",
            "Category and size support for products",
            "Seeder for quick catalog bootstrap"
        ]),
        ("Security Controls", [
            "JWT-based route protection",
            "Admin credentials from environment variables",
            "Separation of public and protected APIs",
            "Recommended next step: rate limiting"
        ]),
        ("Setup Experience", [
            "Env templates added for all apps",
            "Local MongoDB fallback enabled",
            "Clear startup commands documented",
            "One-command product seeding"
        ]),
        ("Operational Readiness", [
            "Consistent endpoint contracts",
            "Structured error responses",
            "Documented local and deployment flow",
            "Known risks and mitigation plan"
        ]),
        ("Quality and Testing", [
            "Manual E2E flow validated",
            "No frontend errors in updated checkout",
            "No backend controller errors in COD flow",
            "Recommended: add API and UI automation"
        ]),
        ("Performance Opportunities", [
            "Add pagination for large catalogs",
            "Optimize image delivery",
            "Apply indexing for frequent queries",
            "Introduce client-side lazy loading"
        ]),
        ("Risks and Mitigation", [
            "Port conflicts -> enforce service ports",
            "Missing env values -> startup validation",
            "Data drift -> stricter admin workflow checks",
            "Scaling challenges -> monitor and profile"
        ]),
        ("Roadmap", [
            "Phase 1: stability and test coverage",
            "Phase 2: promotions and coupon engine",
            "Phase 3: advanced analytics",
            "Phase 4: optional online payment reintroduction"
        ]),
        ("Summary", [
            "Platform is functional and structured",
            "COD-only model is now consistent end-to-end",
            "Documentation and setup are production-oriented",
            "Strong base for future feature expansion"
        ]),
        ("Thank You", [
            "Questions and discussion",
            "Next sprint planning and ownership alignment"
        ]),
    ]

    for title, bullets in slides:
        add_bullets_slide(prs, title, bullets)

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()
    print(f"Created: {DOCX_PATH}")
    print(f"Created: {PPTX_PATH}")
